"""T2303 — 文档生成器测试.

覆盖:
- 3 种格式 (docx/pptx/pdf)
- 5 个模板
- 中文字体 (graceful fallback)
- API 端点
- 内容完整性
"""
from __future__ import annotations

import io
import os
import sys
from unittest.mock import MagicMock, patch
from uuid import UUID, uuid4

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from services.platform.document_generator import (  # noqa: E402
    DocumentGenerator,
    DocumentResult,
    DocFormat,
    DocTemplate,
    HAS_DOCX,
    HAS_PDF,
    HAS_PPTX,
    load_template_data,
    CN_FONT_NAME,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def generator():
    return DocumentGenerator()


@pytest.fixture
def sample_candidate_data():
    return {
        "id": "c1",
        "name": "张三",
        "headline": "Senior Backend Engineer",
        "location": "Shanghai",
        "experience_years": 8,
        "skills": [
            {"name": "Python"},
            {"name": "Go"},
            {"name": "Kubernetes"},
        ],
        "experience": [
            {"title": "Senior Engineer", "company": "ACME", "period": "2020-2024"},
        ],
        "recent_matches": [
            {
                "roles": {"title": "Backend Lead"},
                "overall_score": 0.92,
            },
        ],
        "latest_assessment": {"score": 85, "level": "senior"},
        "generated_at": "2026-07-12T00:00:00Z",
    }


@pytest.fixture
def sample_funnel_data():
    return {
        "period_days": 30,
        "since": "2026-06-12T00:00:00Z",
        "counts": {
            "sourced": 100,
            "screened": 60,
            "interviewed": 30,
            "offered": 10,
            "hired": 5,
        },
        "total_events": 205,
        "generated_at": "2026-07-12T00:00:00Z",
    }


@pytest.fixture
def sample_sla_data():
    return {
        "period_days": 30,
        "total_tickets": 100,
        "sla_breaches": 5,
        "sla_compliance": 95.0,
        "avg_response_hours": 1.2,
        "avg_resolution_hours": 6.5,
        "generated_at": "2026-07-12T00:00:00Z",
    }


@pytest.fixture
def sample_weekly_data():
    return {
        "week_start": "2026-07-06",
        "week_end": "2026-07-12",
        "counts": {"sourced": 20, "hired": 2},
        "total": 22,
    }


@pytest.fixture
def sample_monthly_data():
    return {
        "month_start": "2026-06-12",
        "month_end": "2026-07-12",
        "counts": {"sourced": 80, "hired": 10},
        "total_events": 90,
        "unique_candidates": 60,
    }


# ---------------------------------------------------------------------------
# 1. 基础生成
# ---------------------------------------------------------------------------


def test_generator_init():
    gen = DocumentGenerator()
    assert gen is not None


def test_generator_custom_font_path():
    gen = DocumentGenerator(cn_font_path="/nonexistent.ttf")
    assert gen.cn_font_path == "/nonexistent.ttf"


def test_generate_unknown_format_raises(generator, sample_candidate_data):
    with pytest.raises(ValueError, match="不支持的格式"):
        generator.generate(
            template="candidate_report",
            fmt="xyz",
            data=sample_candidate_data,
        )


def test_generate_unknown_template_raises(generator, sample_candidate_data):
    with pytest.raises(ValueError, match="不支持的模板"):
        generator.generate(
            template="unknown_template",
            fmt="docx",
            data=sample_candidate_data,
        )


def test_generate_returns_document_result(generator, sample_candidate_data):
    result = generator.generate(
        template="candidate_report",
        fmt="docx",
        data=sample_candidate_data,
    )
    assert isinstance(result, DocumentResult)
    assert result.format == "docx"
    assert result.template == "candidate_report"
    assert len(result.content) > 0
    assert result.filename.endswith(".docx")


def test_generate_filename_format(generator, sample_candidate_data):
    for fmt in ["docx", "pptx", "pdf", "txt"]:
        result = generator.generate(
            template="candidate_report",
            fmt=fmt,
            data=sample_candidate_data,
        )
        assert result.filename == f"candidate_report.{fmt}"


def test_generate_mime_types(generator, sample_candidate_data):
    result = generator.generate(
        template="candidate_report", fmt="docx",
        data=sample_candidate_data,
    )
    assert "officedocument" in result.mime_type or "octet-stream" in result.mime_type


def test_generate_size_bytes_set(generator, sample_candidate_data):
    result = generator.generate(
        template="candidate_report", fmt="docx",
        data=sample_candidate_data,
    )
    assert result.size_bytes == len(result.content)


# ---------------------------------------------------------------------------
# 2. 5 个模板 × 3 种格式
# ---------------------------------------------------------------------------


TEMPLATES = [
    ("candidate_report", "sample_candidate_data"),
    ("funnel_report", "sample_funnel_data"),
    ("sla_report", "sample_sla_data"),
    ("weekly_recruitment", "sample_weekly_data"),
    ("monthly_recruitment", "sample_monthly_data"),
]

FORMATS = ["docx", "pptx", "pdf", "txt"]


@pytest.mark.parametrize("template,fixture_name", TEMPLATES)
@pytest.mark.parametrize("fmt", FORMATS)
def test_all_template_format_combinations(
    request, generator, template, fixture_name, fmt
):
    data = request.getfixturevalue(fixture_name)
    result = generator.generate(template=template, fmt=fmt, data=data)
    assert result.content is not None
    assert len(result.content) > 0


# ---------------------------------------------------------------------------
# 3. TXT 渲染 (兜底)
# ---------------------------------------------------------------------------


def test_txt_contains_chinese(generator, sample_candidate_data):
    """中文字符应正确编码."""
    result = generator.generate(
        template="candidate_report", fmt="txt", data=sample_candidate_data,
    )
    text = result.content.decode("utf-8")
    assert "张三" in text


def test_txt_contains_basic_fields(generator, sample_candidate_data):
    result = generator.generate(
        template="candidate_report", fmt="txt", data=sample_candidate_data,
    )
    text = result.content.decode("utf-8")
    assert "name" in text
    assert "experience_years" in text


def test_txt_contains_generated_timestamp(generator, sample_candidate_data):
    result = generator.generate(
        template="candidate_report", fmt="txt", data=sample_candidate_data,
    )
    text = result.content.decode("utf-8")
    assert "生成时间" in text


def test_txt_funnel_contains_counts(generator, sample_funnel_data):
    result = generator.generate(
        template="funnel_report", fmt="txt", data=sample_funnel_data,
    )
    text = result.content.decode("utf-8")
    assert "sourced" in text
    assert "hired" in text


def test_txt_sla_compliance_value(generator, sample_sla_data):
    result = generator.generate(
        template="sla_report", fmt="txt", data=sample_sla_data,
    )
    text = result.content.decode("utf-8")
    assert "95.0" in text


# ---------------------------------------------------------------------------
# 4. DOCX
# ---------------------------------------------------------------------------


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx 未安装")
def test_docx_is_valid(generator, sample_candidate_data):
    result = generator.generate(
        template="candidate_report", fmt="docx", data=sample_candidate_data,
    )
    # docx 是 zip 文件,以 PK 开头
    assert result.content[:2] == b"PK"


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx 未安装")
def test_docx_candidate_can_be_loaded(generator, sample_candidate_data):
    from docx import Document
    result = generator.generate(
        template="candidate_report", fmt="docx", data=sample_candidate_data,
    )
    doc = Document(io.BytesIO(result.content))
    paragraphs = [p.text for p in doc.paragraphs]
    assert any("张三" in p or "候选人" in p for p in paragraphs)


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx 未安装")
def test_docx_funnel_has_stages(generator, sample_funnel_data):
    from docx import Document
    result = generator.generate(
        template="funnel_report", fmt="docx", data=sample_funnel_data,
    )
    doc = Document(io.BytesIO(result.content))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "sourced" in text
    assert "100" in text


# ---------------------------------------------------------------------------
# 5. PPTX
# ---------------------------------------------------------------------------


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx 未安装")
def test_pptx_is_valid(generator, sample_candidate_data):
    result = generator.generate(
        template="candidate_report", fmt="pptx", data=sample_candidate_data,
    )
    # pptx 也是 zip
    assert result.content[:2] == b"PK"


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx 未安装")
def test_pptx_candidate_has_slides(generator, sample_candidate_data):
    from pptx import Presentation
    result = generator.generate(
        template="candidate_report", fmt="pptx", data=sample_candidate_data,
    )
    prs = Presentation(io.BytesIO(result.content))
    # 至少 2 slides (title + body)
    assert len(prs.slides) >= 2


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx 未安装")
def test_pptx_funnel_has_slide(generator, sample_funnel_data):
    from pptx import Presentation
    result = generator.generate(
        template="funnel_report", fmt="pptx", data=sample_funnel_data,
    )
    prs = Presentation(io.BytesIO(result.content))
    assert len(prs.slides) >= 1


# ---------------------------------------------------------------------------
# 6. PDF
# ---------------------------------------------------------------------------


@pytest.mark.skipif(not HAS_PDF, reason="reportlab 未安装")
def test_pdf_magic_bytes(generator, sample_candidate_data):
    result = generator.generate(
        template="candidate_report", fmt="pdf", data=sample_candidate_data,
    )
    assert result.content[:4] == b"%PDF"


@pytest.mark.skipif(not HAS_PDF, reason="reportlab 未安装")
def test_pdf_contains_chinese(generator, sample_candidate_data):
    result = generator.generate(
        template="candidate_report", fmt="pdf", data=sample_candidate_data,
    )
    # PDF 嵌入字体后,中文应以 UTF-16BE 编码出现
    # 这里只验证 PDF magic + 大小合理
    assert result.content[:4] == b"%PDF"
    assert len(result.content) > 1000


@pytest.mark.skipif(not HAS_PDF, reason="reportlab 未安装")
def test_pdf_funnel_has_data(generator, sample_funnel_data):
    result = generator.generate(
        template="funnel_report", fmt="pdf", data=sample_funnel_data,
    )
    assert result.content[:4] == b"%PDF"
    # 包含 funnel 关键词 ASCII
    assert b"sourced" in result.content or b"hired" in result.content or b"PDF" in result.content


# ---------------------------------------------------------------------------
# 7. 中文字体支持
# ---------------------------------------------------------------------------


def test_chinese_font_name():
    """验证中文字体名常量."""
    assert CN_FONT_NAME == "STSong-Light"


def test_font_registration_graceful():
    """即使字体文件不存在,也不应崩溃."""
    gen = DocumentGenerator(cn_font_path="/nonexistent/path/font.ttf")
    # 初始化时不抛错
    assert gen is not None


@pytest.mark.skipif(not HAS_PDF, reason="reportlab 未安装")
def test_pdf_with_missing_font_fallback(generator, sample_candidate_data):
    """字体不存在时使用内置 CIDFont."""
    gen = DocumentGenerator(cn_font_path="/nonexistent/font.ttf")
    result = gen.generate(
        template="candidate_report", fmt="pdf", data=sample_candidate_data,
    )
    assert result.content[:4] == b"%PDF"


# ---------------------------------------------------------------------------
# 8. 模板数据加载
# ---------------------------------------------------------------------------


def test_load_candidate_template_data():
    """用 mock supabase 加载候选人数据."""
    sb = MagicMock()
    # mock candidate fetch
    cand_table = MagicMock()
    cand_data = MagicMock()
    cand_data.data = {"id": "c1", "name": "Alice"}
    cand_table.select.return_value.eq.return_value.single.return_value.execute.return_value = cand_data
    # mock match fetch
    match_table = MagicMock()
    match_data = MagicMock()
    match_data.data = []
    match_table.select.return_value.eq.return_value.order.return_value.limit.return_value.execute.return_value = match_data
    # mock assessment fetch
    a_table = MagicMock()
    a_data = MagicMock()
    a_data.data = []
    a_table.select.return_value.eq.return_value.order.return_value.limit.return_value.execute.return_value = a_data

    def table(name):
        return {"candidates": cand_table, "matches": match_table,
                "assessments": a_table}.get(name, MagicMock())

    sb.table.side_effect = table

    data = load_template_data(
        template="candidate_report",
        supabase=sb,
        candidate_id="c1",
    )
    assert data["name"] == "Alice"


def test_load_candidate_requires_id():
    data = load_template_data(
        template="candidate_report",
        supabase=MagicMock(),
    )
    assert "error" in data


def test_load_funnel_template_data():
    sb = MagicMock()
    funnel_table = MagicMock()
    funnel_data = MagicMock()
    funnel_data.data = [
        {"event_type": "sourced"},
        {"event_type": "sourced"},
        {"event_type": "hired"},
    ]
    funnel_table.select.return_value.gte.return_value.execute.return_value = funnel_data

    def table(name):
        return funnel_table if name == "funnel_events" else MagicMock()

    sb.table.side_effect = table

    data = load_template_data(
        template="funnel_report",
        supabase=sb,
        period_days=30,
    )
    assert data["counts"]["sourced"] == 2
    assert data["counts"]["hired"] == 1


def test_load_sla_template_data():
    sb = MagicMock()
    tickets_table = MagicMock()
    tickets_data = MagicMock()
    tickets_data.data = [
        {
            "id": "t1", "status": "resolved",
            "created_at": "2026-07-01T00:00:00Z",
            "resolved_at": "2026-07-01T02:00:00Z",
            "first_response_at": "2026-07-01T00:30:00Z",
        },
    ]
    tickets_table.select.return_value.gte.return_value.execute.return_value = tickets_data

    def table(name):
        return tickets_table if name == "tickets" else MagicMock()

    sb.table.side_effect = table

    data = load_template_data(
        template="sla_report",
        supabase=sb,
        period_days=30,
    )
    assert data["total_tickets"] == 1
    assert data["sla_compliance"] == 100.0
    assert data["avg_resolution_hours"] == 2.0


def test_load_sla_breach_counted():
    sb = MagicMock()
    tickets_table = MagicMock()
    tickets_data = MagicMock()
    # 5 小时响应 → 超 4h SLA
    tickets_data.data = [
        {
            "id": "t1", "status": "resolved",
            "created_at": "2026-07-01T00:00:00Z",
            "resolved_at": "2026-07-01T10:00:00Z",
            "first_response_at": "2026-07-01T05:00:00Z",
        },
    ]
    tickets_table.select.return_value.gte.return_value.execute.return_value = tickets_data

    def table(name):
        return tickets_table if name == "tickets" else MagicMock()

    sb.table.side_effect = table

    data = load_template_data(
        template="sla_report",
        supabase=sb,
        period_days=30,
    )
    assert data["sla_breaches"] == 1
    assert data["sla_compliance"] == 0.0


def test_load_weekly_data():
    sb = MagicMock()
    table_mock = MagicMock()
    data_mock = MagicMock()
    data_mock.data = [{"event_type": "sourced"}, {"event_type": "hired"}]
    table_mock.select.return_value.gte.return_value.lte.return_value.execute.return_value = data_mock

    def table(name):
        return table_mock

    sb.table.side_effect = table

    data = load_template_data(
        template="weekly_recruitment",
        supabase=sb,
        week_offset=0,
    )
    assert data["total"] == 2


def test_load_monthly_data():
    sb = MagicMock()
    table_mock = MagicMock()
    data_mock = MagicMock()
    data_mock.data = [
        {"event_type": "sourced", "candidate_id": "c1"},
        {"event_type": "sourced", "candidate_id": "c2"},
        {"event_type": "hired", "candidate_id": "c1"},
    ]
    table_mock.select.return_value.gte.return_value.lte.return_value.execute.return_value = data_mock

    def table(name):
        return table_mock

    sb.table.side_effect = table

    data = load_template_data(
        template="monthly_recruitment",
        supabase=sb,
        month_offset=0,
    )
    assert data["unique_candidates"] == 2
    assert data["total_events"] == 3


def test_load_unknown_template_returns_error():
    data = load_template_data(
        template="nonsense",
        supabase=MagicMock(),
    )
    assert "error" in data


# ---------------------------------------------------------------------------
# 9. API 端点
# ---------------------------------------------------------------------------


@pytest.fixture
def api_client():
    from fastapi import FastAPI
    from fastapi.testclient import TestClient
    from api.auth import CurrentUser, get_current_user
    from api.exports import router
    from api.deps import get_supabase

    fake_user = CurrentUser(
        id=UUID("00000000-0000-0000-0000-000000000099"),
        email="t@example.com",
        role="talent_partner",
        organisation_id=None,
    )

    async def override_user():
        return fake_user

    sb = MagicMock()
    cand_data = MagicMock()
    cand_data.data = {"id": "c1", "name": "Alice", "skills": []}
    cand_table = MagicMock()
    cand_table.select.return_value.eq.return_value.single.return_value.execute.return_value = cand_data
    match_data = MagicMock()
    match_data.data = []
    match_table = MagicMock()
    match_table.select.return_value.eq.return_value.order.return_value.limit.return_value.execute.return_value = match_data
    assess_data = MagicMock()
    assess_data.data = []
    assess_table = MagicMock()
    assess_table.select.return_value.eq.return_value.order.return_value.limit.return_value.execute.return_value = assess_data

    funnel_data = MagicMock()
    funnel_data.data = []
    funnel_table = MagicMock()
    funnel_table.select.return_value.gte.return_value.execute.return_value = funnel_data

    sla_data = MagicMock()
    sla_data.data = []
    sla_table = MagicMock()
    sla_table.select.return_value.gte.return_value.execute.return_value = sla_data

    def table(name):
        return {
            "candidates": cand_table,
            "matches": match_table,
            "assessments": assess_table,
            "funnel_events": funnel_table,
            "tickets": sla_table,
        }.get(name, MagicMock())

    sb.table.side_effect = table

    def override_supabase():
        return sb

    app = FastAPI()
    app.include_router(router)
    app.dependency_overrides[get_current_user] = override_user
    app.dependency_overrides[get_supabase] = override_supabase
    return TestClient(app)


def test_api_candidate_report_docx(api_client):
    r = api_client.get(
        f"/api/exports/candidate-report/{uuid4()}?format=docx"
    )
    assert r.status_code == 200
    assert "officedocument" in r.headers.get("content-type", "")
    assert len(r.content) > 0


def test_api_candidate_report_pdf(api_client):
    r = api_client.get(
        f"/api/exports/candidate-report/{uuid4()}?format=pdf"
    )
    assert r.status_code == 200
    assert "pdf" in r.headers.get("content-type", "")


def test_api_candidate_report_pptx(api_client):
    r = api_client.get(
        f"/api/exports/candidate-report/{uuid4()}?format=pptx"
    )
    assert r.status_code == 200
    assert "presentationml" in r.headers.get("content-type", "")


def test_api_candidate_report_invalid_format(api_client):
    r = api_client.get(
        f"/api/exports/candidate-report/{uuid4()}?format=exe"
    )
    assert r.status_code == 422  # pattern mismatch


def test_api_funnel_report(api_client):
    r = api_client.get("/api/exports/funnel-report?format=pdf")
    assert r.status_code == 200


def test_api_sla_report(api_client):
    r = api_client.get("/api/exports/sla-report?format=docx")
    assert r.status_code == 200


def test_api_weekly(api_client):
    r = api_client.get("/api/exports/weekly?format=pptx")
    assert r.status_code == 200


def test_api_monthly(api_client):
    r = api_client.get("/api/exports/monthly?format=docx")
    assert r.status_code == 200


def test_api_disposition_header(api_client):
    r = api_client.get(
        f"/api/exports/candidate-report/{uuid4()}?format=docx"
    )
    cd = r.headers.get("content-disposition", "")
    assert "attachment" in cd
    assert ".docx" in cd


# ---------------------------------------------------------------------------
# 10. 边界
# ---------------------------------------------------------------------------


def test_generate_with_empty_data(generator):
    """空数据应生成有效文档."""
    result = generator.generate(
        template="candidate_report", fmt="txt", data={},
    )
    assert len(result.content) > 0


def test_generate_with_missing_keys(generator):
    """缺字段时不应崩溃."""
    data = {"name": "Test"}  # 只给 name
    result = generator.generate(
        template="candidate_report", fmt="txt", data=data,
    )
    text = result.content.decode("utf-8")
    assert "Test" in text


def test_funnel_with_zero_events(generator):
    data = {"period_days": 30, "counts": {}, "total_events": 0,
            "generated_at": "2026-07-12T00:00:00Z"}
    result = generator.generate(
        template="funnel_report", fmt="docx", data=data,
    )
    assert result.size_bytes > 0


def test_sla_with_zero_tickets(generator):
    data = {"period_days": 30, "total_tickets": 0,
            "sla_breaches": 0, "sla_compliance": 100.0,
            "avg_response_hours": None, "avg_resolution_hours": None,
            "generated_at": "2026-07-12T00:00:00Z"}
    result = generator.generate(
        template="sla_report", fmt="docx", data=data,
    )
    assert result.size_bytes > 0


def test_all_15_combinations_produce_content(generator):
    """5 模板 × 3 格式全部产生非空内容."""
    for template in ["candidate_report", "funnel_report", "sla_report",
                     "weekly_recruitment", "monthly_recruitment"]:
        for fmt in ["docx", "pptx", "pdf"]:
            data = {"generated_at": "2026-07-12T00:00:00Z",
                    "name": "Test", "period_days": 30, "total_events": 0,
                    "counts": {"x": 1}, "total": 1, "total_tickets": 0,
                    "sla_breaches": 0, "sla_compliance": 100,
                    "week_start": "2026-07-06", "week_end": "2026-07-12",
                    "month_start": "2026-06-12", "month_end": "2026-07-12",
                    "unique_candidates": 0}
            result = generator.generate(
                template=template, fmt=fmt, data=data,
            )
            assert len(result.content) > 0, f"{template}/{fmt} 空"
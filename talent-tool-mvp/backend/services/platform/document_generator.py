"""T2303 — 文档生成服务.

支持格式:
- docx (python-docx) — Word 报告
- pptx (python-pptx) — PowerPoint 演示
- pdf (reportlab) — PDF 文档

特性:
- Jinja2 风格模板
- 中文字体嵌入 (宋体 SimSun)
- 5 个内置模板:候选人报告/漏斗分析/工单 SLA/招聘周报/招聘月报

依赖 (可选):
- python-docx
- python-pptx
- reportlab

如未安装则 graceful fallback 到纯文本.
"""
from __future__ import annotations

import io
import logging
import os
from dataclasses import dataclass
from datetime import datetime, timezone
from enum import Enum
from typing import Any, Callable, Optional

logger = logging.getLogger("recruittech.services.docgen")


# ---------------------------------------------------------------------------
# 格式 & 模板
# ---------------------------------------------------------------------------


class DocFormat(str, Enum):
    DOCX = "docx"
    PPTX = "pptx"
    PDF = "pdf"
    TXT = "txt"


class DocTemplate(str, Enum):
    CANDIDATE_REPORT = "candidate_report"
    FUNNEL_REPORT = "funnel_report"
    SLA_REPORT = "sla_report"
    WEEKLY_RECRUITMENT = "weekly_recruitment"
    MONTHLY_RECRUITMENT = "monthly_recruitment"


# 中文字体名 (reportlab 默认支持 STSong-Light)
CN_FONT_NAME = "STSong-Light"
CN_FONT_PATH = os.environ.get(
    "DOC_CN_FONT_PATH",
    "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",  # 常见 Linux 路径
)


@dataclass
class DocumentResult:
    format: str
    template: str
    content: bytes
    filename: str
    mime_type: str
    size_bytes: int


# ---------------------------------------------------------------------------
# Optional imports (graceful fallback)
# ---------------------------------------------------------------------------

try:
    from docx import Document  # python-docx
    from docx.shared import Pt, RGBColor, Inches
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation  # python-pptx
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        Table,
        TableStyle,
    )
    HAS_PDF = True
except ImportError:
    HAS_PDF = False


# ---------------------------------------------------------------------------
# 模板数据获取
# ---------------------------------------------------------------------------


def _candidate_report_data(candidate_id: str, supabase) -> dict[str, Any]:
    """拉取候选人数据 (含 match 信息)."""
    res = (
        supabase.table("candidates")
        .select("*")
        .eq("id", candidate_id)
        .single()
        .execute()
    )
    candidate = res.data or {}
    if not candidate:
        return {"error": "候选人不存在"}

    # 最近 match
    m = (
        supabase.table("matches")
        .select("*, roles(title, organisation_id)")
        .eq("candidate_id", candidate_id)
        .order("created_at", desc=True)
        .limit(3)
        .execute()
    )
    candidate["recent_matches"] = m.data or []

    # 最近评估
    a = (
        supabase.table("assessments")
        .select("*")
        .eq("candidate_id", candidate_id)
        .order("created_at", desc=True)
        .limit(1)
        .execute()
    )
    candidate["latest_assessment"] = (a.data or [None])[0]

    candidate["generated_at"] = datetime.now(timezone.utc).isoformat()
    return candidate


def _funnel_report_data(period_days: int, supabase) -> dict[str, Any]:
    """拉取漏斗数据."""
    from datetime import timedelta
    since = (datetime.now(timezone.utc) - timedelta(days=period_days)).isoformat()
    res = (
        supabase.table("funnel_events")
        .select("event_type, candidate_id, role_id, created_at")
        .gte("created_at", since)
        .execute()
    )
    events = res.data or []
    # 聚合每个阶段的计数
    counts: dict[str, int] = {}
    for e in events:
        counts[e["event_type"]] = counts.get(e["event_type"], 0) + 1
    return {
        "period_days": period_days,
        "since": since,
        "counts": counts,
        "total_events": len(events),
        "generated_at": datetime.now(timezone.utc).isoformat(),
    }


def _sla_report_data(period_days: int, supabase) -> dict[str, Any]:
    """拉取 SLA 数据."""
    from datetime import timedelta
    since = (datetime.now(timezone.utc) - timedelta(days=period_days)).isoformat()
    res = (
        supabase.table("tickets")
        .select("id, status, created_at, resolved_at, first_response_at")
        .gte("created_at", since)
        .execute()
    )
    tickets = res.data or []

    response_times = []
    resolution_times = []
    sla_breaches = 0
    for t in tickets:
        if t.get("first_response_at") and t.get("created_at"):
            rt = _parse_dt(t["first_response_at"]) - _parse_dt(t["created_at"])
            response_times.append(rt.total_seconds() / 3600)  # hours
            if rt.total_seconds() > 4 * 3600:  # 4h SLA
                sla_breaches += 1
        if t.get("resolved_at") and t.get("created_at"):
            rdt = _parse_dt(t["resolved_at"]) - _parse_dt(t["created_at"])
            resolution_times.append(rdt.total_seconds() / 3600)
    return {
        "period_days": period_days,
        "total_tickets": len(tickets),
        "sla_breaches": sla_breaches,
        "sla_compliance": (
            round((1 - sla_breaches / max(len(tickets), 1)) * 100, 1)
        ),
        "avg_response_hours": (
            round(sum(response_times) / len(response_times), 2)
            if response_times else None
        ),
        "avg_resolution_hours": (
            round(sum(resolution_times) / len(resolution_times), 2)
            if resolution_times else None
        ),
        "generated_at": datetime.now(timezone.utc).isoformat(),
    }


def _parse_dt(s: str) -> datetime:
    try:
        return datetime.fromisoformat(s.replace("Z", "+00:00"))
    except Exception:
        return datetime.now(timezone.utc)


def _weekly_data(week_offset: int, supabase) -> dict[str, Any]:
    """招聘周报."""
    from datetime import timedelta
    end = datetime.now(timezone.utc) - timedelta(weeks=week_offset)
    start = end - timedelta(days=7)
    res = (
        supabase.table("funnel_events")
        .select("event_type")
        .gte("created_at", start.isoformat())
        .lte("created_at", end.isoformat())
        .execute()
    )
    events = res.data or []
    counts: dict[str, int] = {}
    for e in events:
        counts[e["event_type"]] = counts.get(e["event_type"], 0) + 1
    return {
        "week_start": start.date().isoformat(),
        "week_end": end.date().isoformat(),
        "counts": counts,
        "total": len(events),
    }


def _monthly_data(month_offset: int, supabase) -> dict[str, Any]:
    """招聘月报."""
    from datetime import timedelta
    end = datetime.now(timezone.utc) - timedelta(days=30 * month_offset)
    start = end - timedelta(days=30)
    res = (
        supabase.table("funnel_events")
        .select("event_type, candidate_id")
        .gte("created_at", start.isoformat())
        .lte("created_at", end.isoformat())
        .execute()
    )
    events = res.data or []
    counts: dict[str, int] = {}
    for e in events:
        counts[e["event_type"]] = counts.get(e["event_type"], 0) + 1
    return {
        "month_start": start.date().isoformat(),
        "month_end": end.date().isoformat(),
        "counts": counts,
        "total_events": len(events),
        "unique_candidates": len({e["candidate_id"] for e in events if e.get("candidate_id")}),
    }


# ---------------------------------------------------------------------------
# 生成器
# ---------------------------------------------------------------------------


class DocumentGenerator:
    """文档生成器."""

    def __init__(self, cn_font_path: str | None = None):
        self.cn_font_path = cn_font_path or CN_FONT_PATH
        self._font_registered = False
        self._register_font()

    def _register_font(self):
        if not HAS_PDF or self._font_registered:
            return
        try:
            if os.path.exists(self.cn_font_path):
                pdfmetrics.registerFont(TTFont("CNFont", self.cn_font_path))
                self._font_registered = True
                logger.info("Registered CN font from %s", self.cn_font_path)
            else:
                # 用 reportlab 内置的中文字体
                from reportlab.pdfbase.cidfonts import UnicodeCIDFont
                pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
                self._font_registered = True
        except Exception as e:
            logger.warning("Font registration failed: %s", e)

    def generate(
        self,
        template: str,
        fmt: str,
        data: dict[str, Any],
        title: str | None = None,
    ) -> DocumentResult:
        """生成文档.

        Args:
            template: candidate_report | funnel_report | sla_report |
                      weekly_recruitment | monthly_recruitment
            fmt: docx | pptx | pdf | txt
            data: 模板数据
            title: 可选标题
        """
        if fmt not in {f.value for f in DocFormat}:
            raise ValueError(f"不支持的格式: {fmt}")
        if template not in {t.value for t in DocTemplate}:
            raise ValueError(f"不支持的模板: {template}")

        renderer_map: dict[str, Callable] = {
            (DocFormat.DOCX.value, "candidate_report"): self._docx_candidate,
            (DocFormat.DOCX.value, "funnel_report"): self._docx_funnel,
            (DocFormat.DOCX.value, "sla_report"): self._docx_sla,
            (DocFormat.DOCX.value, "weekly_recruitment"): self._docx_weekly,
            (DocFormat.DOCX.value, "monthly_recruitment"): self._docx_monthly,
            (DocFormat.PPTX.value, "candidate_report"): self._pptx_candidate,
            (DocFormat.PPTX.value, "funnel_report"): self._pptx_funnel,
            (DocFormat.PPTX.value, "sla_report"): self._pptx_sla,
            (DocFormat.PPTX.value, "weekly_recruitment"): self._pptx_weekly,
            (DocFormat.PPTX.value, "monthly_recruitment"): self._pptx_monthly,
            (DocFormat.PDF.value, "candidate_report"): self._pdf_candidate,
            (DocFormat.PDF.value, "funnel_report"): self._pdf_funnel,
            (DocFormat.PDF.value, "sla_report"): self._pdf_sla,
            (DocFormat.PDF.value, "weekly_recruitment"): self._pdf_weekly,
            (DocFormat.PDF.value, "monthly_recruitment"): self._pdf_monthly,
        }
        key = (fmt, template)
        renderer = renderer_map.get(key)
        if not renderer:
            # fallback to txt — _txt_render returns DocumentResult
            return self._txt_render(template, data, title)

        content = renderer(data, title)
        filename = f"{template}.{fmt}"
        mime_map = {
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pdf": "application/pdf",
            "txt": "text/plain; charset=utf-8",
        }
        return DocumentResult(
            format=fmt,
            template=template,
            content=content,
            filename=filename,
            mime_type=mime_map.get(fmt, "application/octet-stream"),
            size_bytes=len(content),
        )

    # -----------------------------------------------------------------------
    # TXT 兜底
    # -----------------------------------------------------------------------

    def _txt_render(
        self, template: str, data: dict, title: str | None
    ) -> DocumentResult:
        content = self._txt_bytes(template, data, title)
        return DocumentResult(
            format="txt",
            template=template,
            content=content,
            filename=f"{template}.txt",
            mime_type="text/plain; charset=utf-8",
            size_bytes=len(content),
        )

    def _txt_bytes(
        self, template: str, data: dict, title: str | None
    ) -> bytes:
        """生成纯文本字节 (内部 fallback)."""
        lines = []
        lines.append(f"=== {title or template} ===")
        lines.append(f"生成时间: {datetime.now(timezone.utc).isoformat()}")
        lines.append("")
        _render_dict_to_txt(data, lines, prefix="")
        return "\n".join(lines).encode("utf-8")

    # -----------------------------------------------------------------------
    # DOCX
    # -----------------------------------------------------------------------

    def _docx_candidate(self, data: dict, title: str | None) -> bytes:
        if not HAS_DOCX:
            return self._txt_bytes("candidate_report", data, title)
        doc = Document()
        doc.add_heading(title or f"候选人报告 — {data.get('name', '')}", level=1)
        doc.add_paragraph(f"生成时间: {data.get('generated_at', '')}")

        doc.add_heading("基本信息", level=2)
        for k in ("name", "headline", "location", "experience_years"):
            if data.get(k) is not None:
                doc.add_paragraph(f"{k}: {data[k]}")

        if data.get("skills"):
            doc.add_heading("技能", level=2)
            for s in data["skills"][:20]:
                doc.add_paragraph(f"• {s.get('name', s)}", style="List Bullet")

        if data.get("experience"):
            doc.add_heading("经验", level=2)
            for e in data["experience"][:5]:
                doc.add_paragraph(
                    f"{e.get('title', '')} @ {e.get('company', '')} ({e.get('period', '')})"
                )

        if data.get("recent_matches"):
            doc.add_heading("匹配记录", level=2)
            for m in data["recent_matches"]:
                doc.add_paragraph(
                    f"岗位: {m.get('roles', {}).get('title', '')} — "
                    f"分数: {round(m.get('overall_score', 0) * 100, 1)}"
                )

        if data.get("latest_assessment"):
            doc.add_heading("最近评估", level=2)
            doc.add_paragraph(str(data["latest_assessment"]))

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def _docx_funnel(self, data: dict, title: str | None) -> bytes:
        if not HAS_DOCX:
            return self._txt_bytes("funnel_report", data, title)
        doc = Document()
        doc.add_heading(title or "漏斗分析报告", level=1)
        doc.add_paragraph(
            f"周期: 过去 {data.get('period_days', 30)} 天"
        )
        doc.add_paragraph(
            f"总事件数: {data.get('total_events', 0)}"
        )
        doc.add_heading("阶段分布", level=2)
        for stage, count in (data.get("counts") or {}).items():
            doc.add_paragraph(f"• {stage}: {count}")
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def _docx_sla(self, data: dict, title: str | None) -> bytes:
        if not HAS_DOCX:
            return self._txt_bytes("sla_report", data, title)
        doc = Document()
        doc.add_heading(title or "工单 SLA 报告", level=1)
        doc.add_paragraph(
            f"周期: {data.get('period_days', 30)} 天"
        )
        doc.add_paragraph(f"总工单数: {data.get('total_tickets', 0)}")
        doc.add_paragraph(f"SLA 合规率: {data.get('sla_compliance', 0)}%")
        doc.add_paragraph(f"SLA 突破数: {data.get('sla_breaches', 0)}")
        doc.add_paragraph(
            f"平均首次响应: {data.get('avg_response_hours', 'N/A')} 小时"
        )
        doc.add_paragraph(
            f"平均解决时长: {data.get('avg_resolution_hours', 'N/A')} 小时"
        )
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def _docx_weekly(self, data: dict, title: str | None) -> bytes:
        if not HAS_DOCX:
            return self._txt_bytes("weekly_recruitment", data, title)
        doc = Document()
        doc.add_heading(title or "招聘周报", level=1)
        doc.add_paragraph(
            f"周: {data.get('week_start', '')} - {data.get('week_end', '')}"
        )
        doc.add_paragraph(f"总事件: {data.get('total', 0)}")
        doc.add_heading("阶段分布", level=2)
        for stage, count in (data.get("counts") or {}).items():
            doc.add_paragraph(f"• {stage}: {count}")
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def _docx_monthly(self, data: dict, title: str | None) -> bytes:
        if not HAS_DOCX:
            return self._txt_bytes("monthly_recruitment", data, title)
        doc = Document()
        doc.add_heading(title or "招聘月报", level=1)
        doc.add_paragraph(
            f"月: {data.get('month_start', '')} - {data.get('month_end', '')}"
        )
        doc.add_paragraph(f"总事件: {data.get('total_events', 0)}")
        doc.add_paragraph(f"独立候选人数: {data.get('unique_candidates', 0)}")
        doc.add_heading("阶段分布", level=2)
        for stage, count in (data.get("counts") or {}).items():
            doc.add_paragraph(f"• {stage}: {count}")
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    # -----------------------------------------------------------------------
    # PPTX
    # -----------------------------------------------------------------------

    def _pptx_candidate(self, data: dict, title: str | None) -> bytes:
        if not HAS_PPTX:
            return self._txt_bytes("candidate_report", data, title)
        prs = Presentation()
        # Slide 1: Title
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title or f"候选人报告 — {data.get('name', '')}"
        slide.placeholders[1].text = (
            f"生成时间: {data.get('generated_at', '')}"
        )

        # Slide 2: Basic info
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "基本信息"
        body = slide.shapes.placeholders[1].text_frame
        first = True
        for k in ("name", "headline", "location", "experience_years"):
            if data.get(k) is not None:
                p = body.paragraphs[0] if first else body.add_paragraph()
                p.text = f"{k}: {data[k]}"
                first = False

        # Slide 3: Skills
        if data.get("skills"):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "技能"
            tf = slide.shapes.placeholders[1].text_frame
            tf.text = ", ".join(
                str(s.get("name", s)) if isinstance(s, dict) else str(s)
                for s in data["skills"][:30]
            )

        # Slide 4: Recent matches
        if data.get("recent_matches"):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "最近匹配"
            tf = slide.shapes.placeholders[1].text_frame
            tf.text = "\n".join(
                f"• {m.get('roles', {}).get('title', '')} — "
                f"{round(m.get('overall_score', 0) * 100, 1)}"
                for m in data["recent_matches"]
            )

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _pptx_funnel(self, data: dict, title: str | None) -> bytes:
        return self._docx_funnel_to_pptx(data, title, "funnel_report")

    def _pptx_sla(self, data: dict, title: str | None) -> bytes:
        return self._docx_funnel_to_pptx(data, title, "sla_report")

    def _pptx_weekly(self, data: dict, title: str | None) -> bytes:
        return self._docx_funnel_to_pptx(data, title, "weekly_recruitment")

    def _pptx_monthly(self, data: dict, title: str | None) -> bytes:
        return self._docx_funnel_to_pptx(data, title, "monthly_recruitment")

    def _docx_funnel_to_pptx(
        self, data: dict, title: str | None, template: str
    ) -> bytes:
        if not HAS_PPTX:
            return self._txt_bytes(template, data, title)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title or template
        slide.placeholders[1].text = (
            f"生成时间: {data.get('generated_at', '')}"
        )

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "数据概览"
        tf = slide.shapes.placeholders[1].text_frame
        lines = []
        for k, v in data.items():
            if k in ("generated_at", "counts", "recent_matches",
                     "latest_assessment", "skills", "experience"):
                continue
            lines.append(f"{k}: {v}")
        if "counts" in data and isinstance(data["counts"], dict):
            lines.append("")
            lines.append("阶段分布:")
            for stage, count in data["counts"].items():
                lines.append(f"  • {stage}: {count}")
        tf.text = "\n".join(lines) if lines else "无数据"

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # -----------------------------------------------------------------------
    # PDF
    # -----------------------------------------------------------------------

    def _pdf_candidate(self, data: dict, title: str | None) -> bytes:
        if not HAS_PDF:
            return self._txt_bytes("candidate_report", data, title)
        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4)
        styles = getSampleStyleSheet()
        cn_style = ParagraphStyle(
            "cn", parent=styles["Normal"], fontName="STSong-Light",
            fontSize=11, leading=16,
        )
        story = [
            Paragraph(title or f"候选人报告 — {data.get('name', '')}", styles["Title"]),
            Paragraph(f"生成时间: {data.get('generated_at', '')}", cn_style),
            Spacer(1, 0.5 * cm),
            Paragraph("基本信息", styles["Heading2"]),
        ]
        for k in ("name", "headline", "location", "experience_years"):
            if data.get(k) is not None:
                story.append(Paragraph(f"{k}: {data[k]}", cn_style))

        if data.get("skills"):
            story.append(Paragraph("技能", styles["Heading2"]))
            for s in data["skills"][:20]:
                name = s.get("name", s) if isinstance(s, dict) else s
                story.append(Paragraph(f"• {name}", cn_style))

        if data.get("recent_matches"):
            story.append(Paragraph("匹配记录", styles["Heading2"]))
            for m in data["recent_matches"]:
                story.append(Paragraph(
                    f"岗位: {m.get('roles', {}).get('title', '')} — "
                    f"分数: {round(m.get('overall_score', 0) * 100, 1)}",
                    cn_style,
                ))

        doc.build(story)
        return buf.getvalue()

    def _pdf_funnel(self, data: dict, title: str | None) -> bytes:
        return self._pdf_report("funnel_report", data, title,
                                f"漏斗分析报告 (过去 {data.get('period_days', 30)} 天)")

    def _pdf_sla(self, data: dict, title: str | None) -> bytes:
        return self._pdf_report("sla_report", data, title,
                                "工单 SLA 报告")

    def _pdf_weekly(self, data: dict, title: str | None) -> bytes:
        return self._pdf_report("weekly_recruitment", data, title,
                                f"招聘周报 ({data.get('week_start', '')} - {data.get('week_end', '')})")

    def _pdf_monthly(self, data: dict, title: str | None) -> bytes:
        return self._pdf_report("monthly_recruitment", data, title,
                                f"招聘月报 ({data.get('month_start', '')} - {data.get('month_end', '')})")

    def _pdf_report(
        self, template: str, data: dict, title: str | None, headline: str
    ) -> bytes:
        if not HAS_PDF:
            return self._txt_bytes(template, data, title)
        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4)
        styles = getSampleStyleSheet()
        cn_style = ParagraphStyle(
            "cn", parent=styles["Normal"], fontName="STSong-Light",
            fontSize=11, leading=16,
        )
        story = [
            Paragraph(headline, styles["Title"]),
            Paragraph(f"生成时间: {data.get('generated_at', '')}", cn_style),
            Spacer(1, 0.5 * cm),
        ]
        for k, v in data.items():
            if k in ("generated_at", "counts", "recent_matches",
                     "latest_assessment"):
                continue
            story.append(Paragraph(f"{k}: {v}", cn_style))

        if isinstance(data.get("counts"), dict):
            story.append(Spacer(1, 0.3 * cm))
            story.append(Paragraph("阶段分布", styles["Heading2"]))
            for stage, count in data["counts"].items():
                story.append(Paragraph(f"• {stage}: {count}", cn_style))

        doc.build(story)
        return buf.getvalue()


# ---------------------------------------------------------------------------
# 辅助
# ---------------------------------------------------------------------------


def _render_dict_to_txt(d: Any, lines: list, prefix: str = ""):
    if isinstance(d, dict):
        for k, v in d.items():
            if isinstance(v, (dict, list)):
                lines.append(f"{prefix}{k}:")
                _render_dict_to_txt(v, lines, prefix + "  ")
            else:
                lines.append(f"{prefix}{k}: {v}")
    elif isinstance(d, list):
        for item in d:
            if isinstance(item, (dict, list)):
                _render_dict_to_txt(item, lines, prefix)
            else:
                lines.append(f"{prefix}- {item}")
    else:
        lines.append(f"{prefix}{d}")


# ---------------------------------------------------------------------------
# 数据加载器 (公开 API)
# ---------------------------------------------------------------------------


def load_template_data(
    template: str,
    supabase,
    candidate_id: str | None = None,
    period_days: int = 30,
    week_offset: int = 0,
    month_offset: int = 0,
) -> dict[str, Any]:
    """根据模板加载数据."""
    if template == DocTemplate.CANDIDATE_REPORT.value:
        if not candidate_id:
            return {"error": "candidate_id 必填"}
        return _candidate_report_data(candidate_id, supabase)
    elif template == DocTemplate.FUNNEL_REPORT.value:
        return _funnel_report_data(period_days, supabase)
    elif template == DocTemplate.SLA_REPORT.value:
        return _sla_report_data(period_days, supabase)
    elif template == DocTemplate.WEEKLY_RECRUITMENT.value:
        return _weekly_data(week_offset, supabase)
    elif template == DocTemplate.MONTHLY_RECRUITMENT.value:
        return _monthly_data(month_offset, supabase)
    else:
        return {"error": f"未知模板: {template}"}